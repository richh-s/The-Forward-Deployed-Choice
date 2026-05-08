from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# Brand colors
NAVY = RGBColor(0x1A, 0x2E, 0x5A)
TEAL = RGBColor(0x00, 0x8B, 0x8B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF4, 0xF6, 0xF9)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
ACCENT = RGBColor(0xFF, 0x8C, 0x00)

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]  # blank


def add_bg(slide, color=LIGHT_GRAY):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, l, t, w, h, color, transparency=0):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, text, l, t, w, h, font_size=18, bold=False, color=DARK_GRAY,
                 align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_para(tf, text, font_size=14, bold=False, color=DARK_GRAY,
             align=PP_ALIGN.LEFT, italic=False, space_before=0, bullet=False):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    if bullet:
        p.level = 1
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p


def header_bar(slide, title, subtitle=None):
    add_rect(slide, 0, 0, 13.33, 1.35, NAVY)
    add_text_box(slide, title, 0.4, 0.15, 12.5, 0.7,
                 font_size=28, bold=True, color=WHITE)
    if subtitle:
        add_text_box(slide, subtitle, 0.4, 0.82, 12.5, 0.45,
                     font_size=14, color=TEAL, italic=False)


def footer(slide, text="Horizon Services Group | Internal AI Knowledge Assistant | Confidential"):
    add_rect(slide, 0, 7.15, 13.33, 0.35, NAVY)
    add_text_box(slide, text, 0.3, 7.17, 12.5, 0.28,
                 font_size=9, color=WHITE, align=PP_ALIGN.CENTER)


def section_label(slide, text, l, t, w=3.5, h=0.35):
    add_rect(slide, l, t, w, h, TEAL)
    add_text_box(slide, text, l + 0.1, t + 0.04, w - 0.2, h - 0.08,
                 font_size=11, bold=True, color=WHITE)


# ─── SLIDE 1 — Title ─────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, NAVY)

# Decorative accent stripe
add_rect(slide, 0, 0, 0.18, 7.5, TEAL)
add_rect(slide, 0.18, 0, 0.06, 7.5, ACCENT)

# Company name
add_text_box(slide, "HORIZON SERVICES GROUP", 0.6, 1.1, 12, 0.65,
             font_size=15, bold=True, color=TEAL, align=PP_ALIGN.LEFT)

# Main title
tb = slide.shapes.add_textbox(Inches(0.6), Inches(1.75), Inches(11), Inches(1.8))
tb.word_wrap = True
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.LEFT
run = p.add_run()
run.text = "Internal AI Knowledge\nAssistant Proposal"
run.font.size = Pt(40)
run.font.bold = True
run.font.color.rgb = WHITE

# Divider
add_rect(slide, 0.6, 3.65, 5.5, 0.05, TEAL)

# Meta info
add_text_box(slide, "TRP1 – Week 12 Capstone Assessment  |  Part 4A – Final Presentation",
             0.6, 3.82, 11.5, 0.4, font_size=12, color=LIGHT_GRAY)
add_text_box(slide, "Prepared by:  Rachel Samson",
             0.6, 4.35, 6, 0.4, font_size=13, bold=True, color=WHITE)
add_text_box(slide, "Submission Date:  May 8, 2026",
             0.6, 4.8, 6, 0.4, font_size=13, color=LIGHT_GRAY)

# ─── SLIDE 2 — Client Problem ────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, LIGHT_GRAY)
header_bar(slide, "Client Problem & Business Impact", "Why the status quo is unsustainable")
footer(slide)

# Left column
section_label(slide, "INFORMATION SOURCES", 0.4, 1.6, 4.0)
sources = ["Policy documents", "Slack messages", "Manager guidance", "Internal notes"]
tb = slide.shapes.add_textbox(Inches(0.4), Inches(2.05), Inches(4.0), Inches(2.0))
tb.word_wrap = True
tf = tb.text_frame
tf.word_wrap = True
tf.paragraphs[0].runs  # touch first para
for i, s in enumerate(sources):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_before = Pt(4)
    r = p.add_run()
    r.text = f"▸  {s}"
    r.font.size = Pt(13)
    r.font.color.rgb = DARK_GRAY

section_label(slide, "BUSINESS IMPACT", 0.4, 4.2, 4.0)
impacts = ["Repetitive manager/HR questions", "Inconsistent policy guidance",
           "Delayed decision-making", "Compliance & security risks", "Reduced productivity"]
tb2 = slide.shapes.add_textbox(Inches(0.4), Inches(4.65), Inches(4.0), Inches(2.2))
tb2.word_wrap = True
tf2 = tb2.text_frame
tf2.word_wrap = True
for i, imp in enumerate(impacts):
    p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
    p.space_before = Pt(3)
    r = p.add_run()
    r.text = f"▸  {imp}"
    r.font.size = Pt(12)
    r.font.color.rgb = DARK_GRAY

# Right column — example questions cards
section_label(slide, "EXAMPLE EMPLOYEE QUESTIONS", 5.0, 1.6, 7.9)
questions = [
    "\"Can I expense this?\"",
    "\"Who approves remote work equipment?\"",
    "\"Can I use AI tools with client data?\"",
]
for idx, q in enumerate(questions):
    ypos = 2.15 + idx * 1.4
    add_rect(slide, 5.0, ypos, 7.9, 1.1, WHITE)
    add_rect(slide, 5.0, ypos, 0.12, 1.1, TEAL)
    add_text_box(slide, q, 5.25, ypos + 0.2, 7.4, 0.7,
                 font_size=14, italic=True, color=NAVY)

# ─── SLIDE 3 — Proposed Solution ─────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, LIGHT_GRAY)
header_bar(slide, "Proposed AI Assistant", "Scope, capabilities, and key features")
footer(slide)

# Left — solution description
section_label(slide, "THE SOLUTION", 0.4, 1.55, 5.8)
sol_items = [
    "Searches approved company policies",
    "Retrieves relevant information",
    "Generates structured answers",
    "Escalates unclear or risky cases",
]
tb = slide.shapes.add_textbox(Inches(0.4), Inches(2.0), Inches(5.8), Inches(1.9))
tb.word_wrap = True
tf = tb.text_frame
tf.word_wrap = True
for i, item in enumerate(sol_items):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_before = Pt(5)
    r = p.add_run()
    r.text = f"✔  {item}"
    r.font.size = Pt(13)
    r.font.color.rgb = DARK_GRAY

section_label(slide, "MVP SCOPE", 0.4, 4.05, 5.8)
scope = ["HR policies", "Expense reimbursement", "Remote work & equipment",
         "Communication & escalation", "AI tool usage"]
tb2 = slide.shapes.add_textbox(Inches(0.4), Inches(4.5), Inches(5.8), Inches(2.3))
tb2.word_wrap = True
tf2 = tb2.text_frame
tf2.word_wrap = True
for i, s in enumerate(scope):
    p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
    p.space_before = Pt(4)
    r = p.add_run()
    r.text = f"▸  {s}"
    r.font.size = Pt(12)
    r.font.color.rgb = DARK_GRAY

# Right — key features cards
section_label(slide, "KEY FEATURES", 6.5, 1.55, 6.4)
features = [
    ("Source-Backed Answers", "Every response cites the specific policy document"),
    ("Confidence Levels", "Transparency on certainty — escalates when low"),
    ("Escalation Guidance", "Routes unresolved questions to the right person"),
    ("AI Safety Guardrails", "Prevents hallucinations and confidentiality leaks"),
]
for idx, (title, desc) in enumerate(features):
    ypos = 2.1 + idx * 1.22
    add_rect(slide, 6.5, ypos, 6.4, 1.05, WHITE)
    add_rect(slide, 6.5, ypos, 0.12, 1.05, ACCENT)
    add_text_box(slide, title, 6.75, ypos + 0.06, 5.9, 0.38,
                 font_size=13, bold=True, color=NAVY)
    add_text_box(slide, desc, 6.75, ypos + 0.5, 5.9, 0.48,
                 font_size=11, color=DARK_GRAY)

# ─── SLIDE 4 — Architecture ───────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, LIGHT_GRAY)
header_bar(slide, "System Architecture & Workflow", "Components and end-to-end data flow")
footer(slide)

# Components column
section_label(slide, "SYSTEM COMPONENTS", 0.4, 1.55, 5.5)
components = [
    "User Interface (Slack / Web Chat)",
    "Knowledge Base",
    "Retrieval Layer",
    "LLM Layer",
    "Guardrails",
    "Escalation System",
    "Logging & Feedback",
]
tb = slide.shapes.add_textbox(Inches(0.4), Inches(2.0), Inches(5.5), Inches(4.5))
tb.word_wrap = True
tf = tb.text_frame
tf.word_wrap = True
for i, c in enumerate(components):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_before = Pt(5)
    r = p.add_run()
    r.text = f"◈  {c}"
    r.font.size = Pt(13)
    r.font.color.rgb = DARK_GRAY

# Workflow column — numbered steps
section_label(slide, "WORKFLOW", 6.2, 1.55, 6.7)
steps = [
    ("1", "Employee submits question"),
    ("2", "Retrieval layer searches policies"),
    ("3", "System checks risks and conflicts"),
    ("4", "AI generates structured response"),
    ("5", "Escalation triggered if needed"),
    ("6", "Interaction logged for improvement"),
]
arrow_y_positions = []
for idx, (num, step) in enumerate(steps):
    ypos = 2.1 + idx * 0.82
    add_rect(slide, 6.2, ypos, 0.45, 0.55, NAVY)
    add_text_box(slide, num, 6.2, ypos + 0.05, 0.45, 0.45,
                 font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(slide, 6.7, ypos, 6.15, 0.55, WHITE)
    add_text_box(slide, step, 6.8, ypos + 0.1, 5.9, 0.4,
                 font_size=13, color=DARK_GRAY)
    arrow_y_positions.append(ypos)

# ─── SLIDE 5 — Prompt Strategy & Guardrails ──────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, LIGHT_GRAY)
header_bar(slide, "Prompt Strategy & Guardrails", "How the AI reasons safely and accurately")
footer(slide)

# Prompt strategy cards
section_label(slide, "PROMPT STRATEGY", 0.4, 1.55, 6.0)
prompts = [
    ("Retrieval Prompt", "Improves policy document search accuracy by focusing the\nsemantic search on relevant policy sections."),
    ("Answer Generation Prompt", "Creates grounded, employee-facing responses that always\ncite the source policy document."),
    ("Escalation Prompt", "Detects conflicting information, uncertainty, confidential\ndata, and approval-related situations."),
]
for idx, (title, desc) in enumerate(prompts):
    ypos = 2.1 + idx * 1.5
    add_rect(slide, 0.4, ypos, 6.0, 1.3, WHITE)
    add_rect(slide, 0.4, ypos, 0.12, 1.3, TEAL)
    add_text_box(slide, title, 0.65, ypos + 0.08, 5.6, 0.4,
                 font_size=13, bold=True, color=NAVY)
    add_text_box(slide, desc, 0.65, ypos + 0.52, 5.6, 0.7,
                 font_size=11, color=DARK_GRAY)

# Guardrails column
section_label(slide, "GUARDRAILS", 6.7, 1.55, 6.2)
guardrails = [
    "Only answer using approved policies",
    "Prevent hallucinated responses",
    "Detect confidentiality risks",
    "Escalate unclear situations",
    "Do not make final approval decisions",
]
for idx, g in enumerate(guardrails):
    ypos = 2.1 + idx * 0.97
    add_rect(slide, 6.7, ypos, 6.2, 0.75, WHITE)
    add_rect(slide, 6.7, ypos, 0.12, 0.75, ACCENT)
    add_text_box(slide, g, 6.95, ypos + 0.18, 5.8, 0.45,
                 font_size=12, color=DARK_GRAY)

# ─── SLIDE 6 — Example User Journey ──────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, LIGHT_GRAY)
header_bar(slide, "Example User Journey", "End-to-end assistant interaction walkthrough")
footer(slide)

# Question box
add_rect(slide, 0.4, 1.55, 12.53, 0.85, NAVY)
add_text_box(slide,
             '  Employee Question:  "Can I use ChatGPT to summarize confidential client interview notes?"',
             0.4, 1.62, 12.4, 0.65, font_size=13, italic=True, color=WHITE)

# Response cards — 3 across top row
top_cards = [
    ("DIRECT ANSWER", "No. Confidential client information should not be pasted into public AI tools.", TEAL),
    ("EXPLANATION", "The AI Usage Policy allows AI tools for drafting and summarization, but confidential client data must not be shared with public AI systems.", NAVY),
    ("SOURCE REFERENCE", "Policy Note 5 — AI Tool Usage", DARK_GRAY),
]
for idx, (label, body, color) in enumerate(top_cards):
    x = 0.4 + idx * 4.3
    add_rect(slide, x, 2.6, 4.1, 2.1, WHITE)
    add_rect(slide, x, 2.6, 4.1, 0.38, color)
    add_text_box(slide, label, x + 0.1, 2.63, 3.9, 0.3,
                 font_size=10, bold=True, color=WHITE)
    add_text_box(slide, body, x + 0.12, 3.07, 3.85, 1.5,
                 font_size=12, color=DARK_GRAY)

# Bottom row
bot_cards = [
    ("CONFIDENCE LEVEL", "High", TEAL),
    ("ESCALATION NEEDED", "Yes", ACCENT),
    ("SUGGESTED NEXT STEP", "Consult your manager or IT/security team for approved alternatives.", NAVY),
]
for idx, (label, body, color) in enumerate(bot_cards):
    x = 0.4 + idx * 4.3
    add_rect(slide, x, 4.9, 4.1, 1.8, WHITE)
    add_rect(slide, x, 4.9, 4.1, 0.38, color)
    add_text_box(slide, label, x + 0.1, 4.93, 3.9, 0.3,
                 font_size=10, bold=True, color=WHITE)
    add_text_box(slide, body, x + 0.12, 5.37, 3.85, 1.2,
                 font_size=13, bold=(label != "SUGGESTED NEXT STEP"), color=DARK_GRAY)

# ─── SLIDE 7 — Design Decisions, Risks & Next Steps ──────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, LIGHT_GRAY)
header_bar(slide, "Design Decisions, Risks & Next Steps", "Rationale, mitigations, and roadmap")
footer(slide)

# Design decisions
section_label(slide, "KEY DESIGN DECISIONS", 0.4, 1.55, 4.1)
decisions = [
    "Retrieval-based AI reduces hallucinations",
    "Official policies override Slack guidance",
    "Human escalation handles uncertain cases",
]
tb = slide.shapes.add_textbox(Inches(0.4), Inches(2.0), Inches(4.1), Inches(2.0))
tb.word_wrap = True
tf = tb.text_frame
tf.word_wrap = True
for i, d in enumerate(decisions):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_before = Pt(8)
    r = p.add_run()
    r.text = f"✔  {d}"
    r.font.size = Pt(12)
    r.font.color.rgb = DARK_GRAY

# Risks
section_label(slide, "RISKS & MITIGATION", 4.8, 1.55, 4.1)
risks = [
    ("Hallucinated Answers", "Use only retrieved policy documents"),
    ("Conflicting Information", "Prioritize official policies; escalate uncertainty"),
    ("Confidential Data Exposure", "Input filtering and confidentiality warnings"),
]
for idx, (risk, mit) in enumerate(risks):
    ypos = 2.05 + idx * 1.45
    add_rect(slide, 4.8, ypos, 4.1, 1.25, WHITE)
    add_rect(slide, 4.8, ypos, 0.12, 1.25, ACCENT)
    add_text_box(slide, risk, 5.05, ypos + 0.07, 3.7, 0.38,
                 font_size=12, bold=True, color=NAVY)
    add_text_box(slide, mit, 5.05, ypos + 0.52, 3.7, 0.65,
                 font_size=11, color=DARK_GRAY)

# Next steps
section_label(slide, "NEXT STEPS", 9.2, 1.55, 3.9)
steps_ns = [
    ("1", "Conduct policy audit"),
    ("2", "Build MVP retrieval system"),
    ("3", "Pilot with HR, Finance & Ops"),
    ("4", "Improve using feedback & analytics"),
]
for idx, (num, step) in enumerate(steps_ns):
    ypos = 2.05 + idx * 1.2
    add_rect(slide, 9.2, ypos, 0.5, 0.9, TEAL)
    add_text_box(slide, num, 9.2, ypos + 0.15, 0.5, 0.55,
                 font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(slide, 9.75, ypos, 3.3, 0.9, WHITE)
    add_text_box(slide, step, 9.88, ypos + 0.2, 3.1, 0.55,
                 font_size=12, color=DARK_GRAY)

# ─── SLIDE 8 — Conclusion ────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, NAVY)
add_rect(slide, 0, 0, 0.18, 7.5, TEAL)
add_rect(slide, 0.18, 0, 0.06, 7.5, ACCENT)

add_text_box(slide, "FINAL RECOMMENDATION", 0.6, 0.9, 12, 0.5,
             font_size=13, bold=True, color=TEAL)
add_text_box(slide, "Conclusion", 0.6, 1.35, 12, 0.75,
             font_size=34, bold=True, color=WHITE)
add_rect(slide, 0.6, 2.15, 5.5, 0.05, TEAL)

# Value proposition
add_text_box(slide,
             "The proposed Internal AI Knowledge Assistant is a practical and scalable\n"
             "solution that improves productivity, consistency, operational efficiency,\n"
             "and compliance safety.",
             0.6, 2.3, 12.2, 1.1, font_size=14, color=LIGHT_GRAY)

# Pillars
pillars = [
    ("Retrieval-Based AI", "Grounded answers from approved policies"),
    ("Strong Guardrails", "Prevents hallucinations and confidentiality risks"),
    ("Source-Backed Answers", "Every response cites the policy document"),
    ("Human Escalation", "Uncertain cases routed to the right person"),
]
for idx, (title, desc) in enumerate(pillars):
    x = 0.5 + idx * 3.2
    add_rect(slide, x, 3.6, 3.0, 1.7, RGBColor(0x22, 0x3A, 0x6E))
    add_rect(slide, x, 3.6, 3.0, 0.12, TEAL)
    add_text_box(slide, title, x + 0.15, 3.78, 2.7, 0.45,
                 font_size=12, bold=True, color=WHITE)
    add_text_box(slide, desc, x + 0.15, 4.28, 2.7, 0.9,
                 font_size=11, color=LIGHT_GRAY)

add_rect(slide, 0.6, 5.55, 11.7, 0.05, TEAL)
add_text_box(slide,
             "Horizon Services Group can safely improve internal knowledge access while reducing operational risks.",
             0.6, 5.7, 12.1, 0.55, font_size=13, italic=True, color=TEAL)

add_text_box(slide, "Thank You", 0.6, 6.35, 6, 0.65,
             font_size=22, bold=True, color=WHITE)
add_text_box(slide, "Rachel Samson  |  May 8, 2026", 0.6, 6.9, 6, 0.4,
             font_size=11, color=LIGHT_GRAY)

# ─── Save ─────────────────────────────────────────────────────────────────────
out_path = "/Users/aman/Desktop/projects/10academy/The-Forward-Deployed-Choice/Horizon_Services_AI_Assistant_Proposal.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
