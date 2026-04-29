"""
Build the Tenacious Consulting demo presentation.
Run: python scripts/build_presentation.py
Output: Tenacious_Consulting_Presentation.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colours ────────────────────────────────────────────────────────────────
BG       = RGBColor(0x0b, 0x0e, 0x18)   # near-black
SURFACE  = RGBColor(0x13, 0x17, 0x26)   # card bg
BLUE     = RGBColor(0x3b, 0x82, 0xf6)   # primary blue
GREEN    = RGBColor(0x10, 0xb9, 0x81)   # pass green
AMBER    = RGBColor(0xf5, 0x9e, 0x0b)   # warning
RED      = RGBColor(0xef, 0x44, 0x44)   # fail
WHITE    = RGBColor(0xff, 0xff, 0xff)
MUTED    = RGBColor(0x64, 0x74, 0x8b)
LAVENDER = RGBColor(0x81, 0x8c, 0xf8)

W = Inches(13.33)   # widescreen width
H = Inches(7.5)     # widescreen height


# ── Helpers ────────────────────────────────────────────────────────────────
def add_slide(prs, layout_idx=6):
    layout = prs.slide_layouts[layout_idx]
    return prs.slides.add_slide(layout)


def fill_bg(slide, colour=BG):
    from pptx.util import Emu
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = colour


def box(slide, left, top, width, height,
        bg=None, border=None, radius=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.line.fill.background() if border is None else None
    if border:
        shape.line.color.rgb = border
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    if bg:
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg
    else:
        shape.fill.background()
    return shape


def txt(slide, text, left, top, width, height,
        size=18, bold=False, colour=WHITE,
        align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = colour
    return txb


def multiline(slide, lines, left, top, width, height,
              size=14, colour=WHITE, bold=False, spacing=1.15):
    """lines: list of (text, colour, bold, size) or just str"""
    from pptx.util import Pt as _Pt
    from pptx.oxml.ns import qn
    from lxml import etree

    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True

    first = True
    for item in lines:
        if isinstance(item, str):
            t, c, b, s = item, colour, bold, size
        else:
            t = item[0]
            c = item[1] if len(item) > 1 else colour
            b = item[2] if len(item) > 2 else bold
            s = item[3] if len(item) > 3 else size

        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = t
        run.font.size = _Pt(s)
        run.font.bold = b
        run.font.color.rgb = c
    return txb


def tag_pill(slide, text, left, top, colour=BLUE):
    bg_col = RGBColor(
        min(colour[0] + 30, 255) if False else colour[0] // 4,
        colour[1] // 4,
        colour[2] // 4,
    )
    b = box(slide, left, top, 1.4, 0.28, bg=bg_col, border=colour)
    txt(slide, text, left + 0.05, top + 0.02, 1.3, 0.26,
        size=10, bold=True, colour=colour, align=PP_ALIGN.CENTER)
    return b


def divider(slide, top, colour=MUTED):
    shape = slide.shapes.add_shape(
        1, Inches(0.5), Inches(top), Inches(12.33), Inches(0.01)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = colour
    shape.line.fill.background()


def slide_title(slide, title, subtitle=None, tag=None, tag_colour=BLUE):
    fill_bg(slide)
    # accent bar left
    b = box(slide, 0, 0, 0.06, 7.5, bg=BLUE)
    txt(slide, title, 0.3, 0.22, 12.5, 0.7,
        size=32, bold=True, colour=WHITE)
    if subtitle:
        txt(slide, subtitle, 0.3, 0.95, 12.5, 0.45,
            size=16, colour=MUTED)
    if tag:
        tag_pill(slide, tag, 0.3, 0.22, tag_colour)
        txt(slide, title, 0.3, 0.50, 12.5, 0.65,
            size=32, bold=True, colour=WHITE)


# ── BUILD SLIDES ────────────────────────────────────────────────────────────
def build():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    # ── 1. TITLE ──────────────────────────────────────────────────────────
    sl = add_slide(prs)
    fill_bg(sl)
    box(sl, 0, 0, 0.08, 7.5, bg=BLUE)
    box(sl, 0.08, 2.8, 13.25, 0.06, bg=BLUE)

    txt(sl, "Tenacious Consulting", 0.4, 1.2, 12.5, 0.9,
        size=48, bold=True, colour=WHITE)
    txt(sl, "Conversion Engine", 0.4, 2.05, 12.5, 0.8,
        size=48, bold=True, colour=BLUE)
    txt(sl, "Forward-Deployed Challenge  ·  April 2026  ·  Rahel Samson",
        0.4, 3.0, 12.5, 0.5, size=16, colour=MUTED)
    txt(sl, "Signal-grounded outreach  →  Email  →  SMS  →  Booking  →  HubSpot",
        0.4, 3.55, 12.5, 0.5, size=14, colour=LAVENDER)

    # ── 2. PROBLEM ────────────────────────────────────────────────────────
    sl = add_slide(prs)
    fill_bg(sl)
    box(sl, 0, 0, 0.08, 7.5, bg=AMBER)
    txt(sl, "The Problem", 0.4, 0.25, 12, 0.6, size=34, bold=True)
    divider(sl, 1.0)

    problems = [
        ("Most outbound fails because it is generic.",          WHITE,  True,  20),
        ("A cold email that doesn't reference what's actually", MUTED,  False, 17),
        ("happening at the company gets ignored.",              MUTED,  False, 17),
        ("",                                                    WHITE,  False, 10),
        ("Tenacious needed a system that:",                     WHITE,  True,  20),
        ("  →  Reads public signals before writing a word",     LAVENDER, False, 17),
        ("  →  Decides whether and how to pitch",               LAVENDER, False, 17),
        ("  →  Handles the full pipeline automatically",        LAVENDER, False, 17),
        ("  →  Logs every claim so graders can verify it",      LAVENDER, False, 17),
    ]
    multiline(sl, problems, 0.4, 1.15, 12.2, 5.8)

    # ── 3. ARCHITECTURE ───────────────────────────────────────────────────
    sl = add_slide(prs)
    fill_bg(sl)
    box(sl, 0, 0, 0.08, 7.5, bg=LAVENDER)
    txt(sl, "Architecture", 0.4, 0.25, 12, 0.6, size=34, bold=True)
    divider(sl, 1.0)

    steps = [
        ("Public Data  (Crunchbase · Wellfound · layoffs.fyi · press)",  MUTED,    False, 13),
        ("          ↓",                                                    MUTED,    False, 13),
        ("Enrichment Pipeline  →  6 Signals  →  hiring_signal_brief",     LAVENDER, True,  14),
        ("          ↓",                                                    MUTED,    False, 13),
        ("Email Agent  (confidence ≥ 0.70 → Assertion  |  < 0.70 → Inquiry)", BLUE, True, 14),
        ("          ↓",                                                    MUTED,    False, 13),
        ("Resend API  →  Prospect Inbox",                                  WHITE,    False, 13),
        ("          ↓  prospect replies",                                  MUTED,    False, 13),
        ("Webhook Server (Render)  →  Intent Classification",              WHITE,    False, 13),
        ("          ↓  WARM",                                              GREEN,    False, 13),
        ("Africa's Talking SMS  →  Cal.com Booking  →  HubSpot Write-back", GREEN,  True,  14),
        ("          ↓",                                                    MUTED,    False, 13),
        ("Langfuse Traces  →  evidence_graph.json  →  invoice_summary.json", AMBER, False, 13),
    ]
    multiline(sl, steps, 0.4, 1.1, 12.5, 6.2, spacing=1.3)

    # ── 4. SIX SIGNALS ────────────────────────────────────────────────────
    sl = add_slide(prs)
    fill_bg(sl)
    box(sl, 0, 0, 0.08, 7.5, bg=BLUE)
    txt(sl, "Card 1 — Six Enrichment Signals", 0.4, 0.25, 12, 0.6,
        size=30, bold=True)
    tag_pill(sl, "Enrichment", 0.4, 0.9, BLUE)
    txt(sl, "NovaPay Technologies", 1.9, 0.88, 10, 0.35, size=14, colour=MUTED)
    divider(sl, 1.3)

    signals = [
        [("Signal 1", BLUE, True, 13),   ("  Funding Event",       WHITE, False, 13), ("  Series B · $16M · Feb 2026 · 67 days ago",          MUTED, False, 12)],
        [("Signal 2", BLUE, True, 13),   ("  Job-Post Velocity",   WHITE, False, 13), ("  9 open engineering roles · +9 delta over 60d",       MUTED, False, 12)],
        [("Signal 3", BLUE, True, 13),   ("  Layoff Event",        WHITE, False, 13), ("  No layoff detected  [HIGH confidence]",              GREEN, False, 12)],
        [("Signal 4", BLUE, True, 13),   ("  Leadership Change",   WHITE, False, 13), ("  New VP Engineering · started March 16 · 38d ago",   AMBER, False, 12)],
        [("Signal 5", BLUE, True, 13),   ("  AI Maturity Score",   WHITE, False, 13), ("  Score 2/3 · ML roles open · Head of AI on team",    MUTED, False, 12)],
        [("Signal 6", GREEN, True, 13),  ("  ICP Segment",         WHITE, False, 13), ("  Segment 3 — Engineering Leadership Transition  ← derived, never fetched", LAVENDER, False, 12)],
    ]

    y = 1.4
    for row in signals:
        x = 0.4
        for part in row:
            w = [1.1, 2.0, 8.5][row.index(part)]
            txt(sl, part[0], x, y, w, 0.35, size=part[3], bold=part[2], colour=part[1])
            x += w
        y += 0.42

    # competitor gap
    divider(sl, 4.15)
    txt(sl, "Competitor Gap", 0.4, 4.25, 3, 0.35, size=13, bold=True, colour=AMBER)
    comps = [
        ("Stripe",  "AI Maturity 3/3", GREEN),
        ("Plaid",   "AI Maturity 3/3", GREEN),
        ("Brex",    "AI Maturity 2/3", AMBER),
        ("NovaPay", "AI Maturity 2/3  ← prospect", BLUE),
    ]
    cx = 0.4
    for name, score, col in comps:
        box(sl, cx, 4.65, 2.9, 0.7, bg=SURFACE, border=col)
        txt(sl, name,  cx+0.12, 4.68, 2.7, 0.3, size=13, bold=True, colour=col)
        txt(sl, score, cx+0.12, 4.95, 2.7, 0.3, size=11, colour=MUTED)
        cx += 3.1

    # ── 5. LIVE EMAIL ─────────────────────────────────────────────────────
    sl = add_slide(prs)
    fill_bg(sl)
    box(sl, 0, 0, 0.08, 7.5, bg=LAVENDER)
    txt(sl, "Card 2 — Live Email End-to-End", 0.4, 0.25, 12, 0.6,
        size=30, bold=True)
    tag_pill(sl, "Core Pipeline", 0.4, 0.9, LAVENDER)
    divider(sl, 1.3)

    # left panel — how it works
    txt(sl, "How it works", 0.4, 1.4, 5.5, 0.4, size=14, bold=True, colour=MUTED)
    logic = [
        ("Avg signal confidence = 0.75", WHITE, True, 14),
        ("≥ 0.70  →  ASSERTION MODE", GREEN, True, 14),
        ("< 0.70  →  Inquiry Mode", AMBER, False, 13),
        ("", WHITE, False, 8),
        ("Agent states facts, not guesses.", WHITE, False, 13),
        ("Sent via Resend API.", MUTED, False, 13),
        ("", WHITE, False, 8),
        ("Trace ID:  908c8ce3", LAVENDER, False, 12),
        ("Contact:   HubSpot #477559194332", LAVENDER, False, 12),
        ("Cost:      $0.03153", GREEN, False, 12),
    ]
    multiline(sl, logic, 0.4, 1.85, 5.5, 4.5)

    # right panel — actual email
    box(sl, 6.2, 1.35, 6.8, 5.5, bg=SURFACE,
        border=RGBColor(0x25, 0x2d, 0x45))
    txt(sl, "Actual email received", 6.35, 1.45, 6.5, 0.35,
        size=11, colour=MUTED, bold=True)
    txt(sl, "Subject: Request: Aligning Engineering Strategies Post-Transition",
        6.35, 1.82, 6.5, 0.4, size=11, bold=True, colour=WHITE)
    divider(sl, 2.3)
    body = (
        "Hi NovaPay team,\n\n"
        "Congratulations on your recent leadership transition\n"
        "and Series B funding! With 9 open engineering roles,\n"
        "it seems like a pivotal time for vendor reassessment.\n\n"
        "Peers like Stripe have invested in dedicated ML\n"
        "platform teams — an opportunity for NovaPay to\n"
        "enhance its AI capabilities.\n\n"
        "Can we schedule a brief call to discuss your roadmap?\n\n"
        "Best,\nAlex Chen\nSenior Engagement Manager\nTenacious Intelligence Corporation"
    )
    txt(sl, body, 6.35, 2.35, 6.5, 4.4, size=11, colour=WHITE)

    # ── 6. HUBSPOT ────────────────────────────────────────────────────────
    sl = add_slide(prs)
    fill_bg(sl)
    box(sl, 0, 0, 0.08, 7.5, bg=RGBColor(0x8b, 0x5c, 0xf6))
    txt(sl, "Card 3 — HubSpot Contact Populating", 0.4, 0.25, 12, 0.6,
        size=30, bold=True)
    tag_pill(sl, "CRM", 0.4, 0.9, RGBColor(0x8b, 0x5c, 0xf6))
    divider(sl, 1.3)

    fields = [
        ("firstname / lastname",       "Jordan  Kim"),
        ("email",                      "rahelsamson953@gmail.com"),
        ("company",                    "NovaPay Technologies"),
        ("funding_round_type",         "Series B"),
        ("funding_days_ago",           "67"),
        ("funding_amount_usd",         "16000000"),
        ("open_engineering_roles",     "9"),
        ("layoff_event_present",       "False"),
        ("leadership_change_role",     "VP Engineering"),
        ("ai_maturity_score",          "2"),
        ("icp_segment",                "segment_3_leadership_transition"),
        ("icp_conflict_flag",          "True"),
        ("meeting_booked",             "false  →  true  (after booking)"),
    ]
    y = 1.4
    for i, (field, value) in enumerate(fields):
        bg_c = RGBColor(0x1c, 0x21, 0x38) if i % 2 == 0 else SURFACE
        box(sl, 0.4, y, 12.3, 0.33, bg=bg_c)
        txt(sl, field, 0.55, y+0.04, 4.5, 0.28, size=11, colour=MUTED)
        txt(sl, value, 5.2,  y+0.04, 7.3, 0.28, size=11, colour=WHITE)
        y += 0.34

    txt(sl, "⚠  No CFPB · No compliance · No regulatory fields anywhere",
        0.4, 6.9, 12, 0.4, size=12, bold=True, colour=AMBER)

    # ── 7. SMS HANDOFF ────────────────────────────────────────────────────
    sl = add_slide(prs)
    fill_bg(sl)
    box(sl, 0, 0, 0.08, 7.5, bg=GREEN)
    txt(sl, "Card 4 — Email-to-SMS Channel Handoff", 0.4, 0.25, 12, 0.6,
        size=30, bold=True)
    tag_pill(sl, "Multi-Channel", 0.4, 0.9, GREEN)
    divider(sl, 1.3)

    steps_sms = [
        ("STEP 1", "Prospect replies:", BLUE),
        ("",       '"Hi — yes, this sounds interesting. Happy to jump on a quick call."', MUTED),
        ("STEP 2", "Intent classification  →  WARM", GREEN),
        ("",       "Keyword match: 'happy to', 'call'", MUTED),
        ("STEP 3", "Africa's Talking fires SMS  →  +251963055269", GREEN),
        ("",       '"Hi Jordan — Book a 15-min call: cal.com/rahel-samson-tmtjxt/15min"', MUTED),
        ("",       "Delivered · 1/1 recipients · Cost $0.02", MUTED),
        ("STEP 4", "Booking confirmed — CAL-E7067FBE · April 29 10:00 UTC", GREEN),
        ("STEP 5", "HubSpot updated — meeting_booked = true · status 200", GREEN),
    ]
    y = 1.45
    for label, content, col in steps_sms:
        if label:
            txt(sl, label, 0.4, y, 1.2, 0.35, size=12, bold=True, colour=col)
            txt(sl, content, 1.7, y, 11, 0.35, size=13, colour=WHITE)
        else:
            txt(sl, content, 1.7, y, 11, 0.35, size=12, colour=MUTED)
        y += 0.44

    txt(sl, "Zero humans in the loop — email reply to booked call fully automated.",
        0.4, 6.85, 12, 0.4, size=13, bold=True, colour=GREEN)

    # ── 8. SAFETY GATE ────────────────────────────────────────────────────
    sl = add_slide(prs)
    fill_bg(sl)
    box(sl, 0, 0, 0.08, 7.5, bg=AMBER)
    txt(sl, "Card 5 — Agent Refuses to Over-Claim", 0.4, 0.25, 12, 0.6,
        size=30, bold=True)
    tag_pill(sl, "Safety Gate", 0.4, 0.9, AMBER)
    divider(sl, 1.3)

    # input box
    box(sl, 0.4, 1.4, 5.8, 2.8, bg=SURFACE,
        border=RGBColor(0x25, 0x2d, 0x45))
    txt(sl, "Input signals", 0.6, 1.5, 5.4, 0.35, size=12, bold=True, colour=MUTED)
    inp = [
        ("Funding:    Series A · $5M · 60 days ago  [HIGH]", WHITE,  False, 13),
        ("Open roles: 2 engineering roles only", AMBER,  True,  13),
        ("Layoff:     None", MUTED,  False, 13),
        ("Leadership: No change", MUTED,  False, 13),
        ("AI maturity: 1/3", MUTED,  False, 13),
    ]
    multiline(sl, inp, 0.6, 1.9, 5.5, 2.1)

    # arrow
    txt(sl, "→", 6.4, 2.6, 0.6, 0.5, size=30, bold=True, colour=AMBER, align=PP_ALIGN.CENTER)

    # output box
    box(sl, 7.1, 1.4, 5.8, 2.8, bg=SURFACE,
        border=RGBColor(0x25, 0x2d, 0x45))
    txt(sl, "Classification", 7.3, 1.5, 5.4, 0.35, size=12, bold=True, colour=MUTED)
    out = [
        ("ABSTAIN", RED, True, 22),
        ("", WHITE, False, 8),
        ("Fresh funding detected but only 2 open", WHITE, False, 13),
        ("engineering roles. Requires ≥ 5 for", WHITE, False, 13),
        ("Segment 1 qualification.", WHITE, False, 13),
    ]
    multiline(sl, out, 7.3, 1.9, 5.5, 2.1)

    txt(sl, "✓  PASS — Agent correctly abstained. No over-claiming.",
        0.4, 4.55, 12, 0.45, size=14, bold=True, colour=GREEN)

    txt(sl, "Why this matters:", 0.4, 5.15, 12, 0.35, size=13, bold=True, colour=WHITE)
    txt(sl,
        "Pitching 'scale up' to a company with 2 open roles destroys credibility immediately.\n"
        "The classifier stops the agent before the email is ever composed.",
        0.4, 5.55, 12, 0.7, size=13, colour=MUTED)

    # ── 9. SEGMENT 2 ──────────────────────────────────────────────────────
    sl = add_slide(prs)
    fill_bg(sl)
    box(sl, 0, 0, 0.08, 7.5, bg=BLUE)
    txt(sl, "Card 6 — Segment 2 Routing", 0.4, 0.25, 12, 0.6, size=30, bold=True)
    tag_pill(sl, "Enrichment", 0.4, 0.9, BLUE)
    txt(sl, "Post-Layoff + Funded — Monte Carlo",
        1.95, 0.88, 10, 0.35, size=14, colour=MUTED)
    divider(sl, 1.3)

    # two signal boxes
    box(sl, 0.4, 1.4, 5.8, 1.7, bg=SURFACE, border=GREEN)
    txt(sl, "Signal 1 — Funding  (Crunchbase ODM)", 0.6, 1.5, 5.4, 0.35,
        size=12, bold=True, colour=GREEN)
    txt(sl, "Series D  ·  $25M  ·  109 days ago  ·  [HIGH confidence]",
        0.6, 1.88, 5.4, 0.6, size=13, colour=WHITE)

    box(sl, 6.7, 1.4, 5.8, 1.7, bg=SURFACE, border=RED)
    txt(sl, "Signal 3 — Layoff  (layoffs.fyi CSV)", 6.9, 1.5, 5.4, 0.35,
        size=12, bold=True, colour=RED)
    txt(sl, "30% of headcount  ·  32 days ago  ·  [HIGH confidence]",
        6.9, 1.88, 5.4, 0.6, size=13, colour=WHITE)

    # comparison
    divider(sl, 3.25)
    txt(sl, "Naive classifier:", 0.4, 3.35, 5.5, 0.35, size=13, bold=True, colour=RED)
    txt(sl, '"Fresh Series D funding  →  Segment 1 — pitch growth, scale up"',
        0.4, 3.72, 12.5, 0.4, size=13, colour=RED)

    txt(sl, "Our classifier:", 0.4, 4.25, 5.5, 0.35, size=13, bold=True, colour=GREEN)
    txt(sl, "Layoff 32d ago overrides funding signal  →  Segment 2 — cost restructuring pitch",
        0.4, 4.62, 12.5, 0.4, size=13, colour=GREEN)
    txt(sl, '"Preserve your AI delivery capacity while reshaping your cost structure."',
        0.4, 5.0, 12.5, 0.4, size=13, bold=True, colour=LAVENDER)

    txt(sl, "Conflict flag raised  ·  Human review triggered before outreach goes out  ·  Confidence 1.00  ·  PASS",
        0.4, 5.6, 12.5, 0.45, size=12, colour=AMBER)

    # ── 10. TAU2 BENCH ────────────────────────────────────────────────────
    sl = add_slide(prs)
    fill_bg(sl)
    box(sl, 0, 0, 0.08, 7.5, bg=BLUE)
    txt(sl, "Card 7 — τ²-Bench Baseline", 0.4, 0.25, 12, 0.6, size=30, bold=True)
    tag_pill(sl, "Evaluation", 0.4, 0.9, BLUE)
    divider(sl, 1.3)

    # big stats
    stats = [
        ("72.67%", "pass@1", GREEN),
        ("150",    "simulations", BLUE),
        ("$0.0199","cost / sim",  AMBER),
        ("105.9s", "p50 latency", LAVENDER),
    ]
    sx = 0.4
    for val, lbl, col in stats:
        box(sl, sx, 1.4, 2.9, 1.4, bg=SURFACE,
            border=RGBColor(0x25, 0x2d, 0x45))
        txt(sl, val, sx+0.15, 1.5, 2.6, 0.75, size=30, bold=True, colour=col)
        txt(sl, lbl, sx+0.15, 2.2, 2.6, 0.4,  size=13, colour=MUTED)
        sx += 3.1

    txt(sl, "95% CI:  [65.04%  —  79.17%]     Model: Qwen3-Next-80B     Domain: retail     30 tasks × 5 trials",
        0.4, 2.95, 12.5, 0.4, size=12, colour=MUTED)
    divider(sl, 3.45)

    # sample traces
    txt(sl, "Sample traces from eval/trace_log.jsonl", 0.4, 3.55, 12, 0.35,
        size=13, bold=True, colour=WHITE)
    traces = [
        ("task  1", "PASS", "reward=1.0", "$0.0175", "106.8s"),
        ("task  2", "PASS", "reward=1.0", "$0.0288", "177.8s"),
        ("task  7", "PASS", "reward=1.0", "$0.0167", "102.2s"),
        ("task 11", "FAIL", "reward=0.0", "$0.0133", "82.6s"),
        ("task 15", "PASS", "reward=1.0", "$0.0107", "57.2s"),
    ]
    y = 3.98
    for task, status, reward, cost, dur in traces:
        sc = GREEN if status == "PASS" else RED
        txt(sl, task,   0.4,  y, 1.2, 0.32, size=12, colour=MUTED)
        txt(sl, status, 1.7,  y, 0.9, 0.32, size=12, bold=True, colour=sc)
        txt(sl, reward, 2.7,  y, 1.8, 0.32, size=12, colour=WHITE)
        txt(sl, cost,   4.6,  y, 1.5, 0.32, size=12, colour=AMBER)
        txt(sl, dur,    6.2,  y, 1.5, 0.32, size=12, colour=MUTED)
        y += 0.36

    # ── 11. PROBES ────────────────────────────────────────────────────────
    sl = add_slide(prs)
    fill_bg(sl)
    box(sl, 0, 0, 0.08, 7.5, bg=AMBER)
    txt(sl, "Card 8 — Probe Library → Concrete Fix", 0.4, 0.25, 12, 0.6,
        size=30, bold=True)
    tag_pill(sl, "Safety", 0.4, 0.9, AMBER)
    divider(sl, 1.3)

    # failure box
    box(sl, 0.4, 1.4, 12.3, 1.2, bg=RGBColor(0x2d, 0x1a, 0x0a),
        border=RED)
    txt(sl, "Highest-ROI Failure: bench_over_commitment  (P-009)",
        0.6, 1.5, 12, 0.4, size=14, bold=True, colour=RED)
    txt(sl, "Baseline trigger rate: 100%   ·   Business cost per occurrence: $18,000 expected pipeline value",
        0.6, 1.88, 12, 0.4, size=13, colour=AMBER)

    # before / after
    box(sl, 0.4, 2.8, 5.8, 1.8, bg=SURFACE, border=RED)
    txt(sl, "BEFORE (baseline)", 0.6, 2.9, 5.4, 0.35, size=12, bold=True, colour=RED)
    txt(sl,
        "Agent promises capacity it doesn't have.\nEvery probe triggered it — 100% failure rate.\nDiscovery call fails. Deal lost.",
        0.6, 3.28, 5.4, 1.2, size=12, colour=WHITE)

    txt(sl, "→", 6.45, 3.4, 0.6, 0.5, size=28, bold=True,
        colour=GREEN, align=PP_ALIGN.CENTER)

    box(sl, 7.1, 2.8, 5.8, 1.8, bg=SURFACE, border=GREEN)
    txt(sl, "AFTER (confidence-gated agent)", 7.3, 2.9, 5.4, 0.35,
        size=12, bold=True, colour=GREEN)
    txt(sl,
        "Threshold 0.70 — below it, routes to human.\nTrigger rate drops to 0%.\np-value = 5.4 × 10⁻⁶",
        7.3, 3.28, 5.4, 1.2, size=12, colour=WHITE)

    # ablation table
    divider(sl, 4.8)
    txt(sl, "Ablation Results", 0.4, 4.9, 12, 0.35, size=13, bold=True, colour=WHITE)
    ablations = [
        ("Baseline",           "95.0%", "No confidence gating",              MUTED),
        ("Mechanism v1",       "95.0%", "Confidence gating + ICP abstention", BLUE),
        ("Mechanism v2 strict","100.0%","Stricter thresholds",                GREEN),
    ]
    y = 5.32
    for name, score, desc, col in ablations:
        txt(sl, name,  0.4, y, 3.0, 0.32, size=12, colour=col, bold=True)
        txt(sl, score, 3.5, y, 1.2, 0.32, size=12, colour=GREEN, bold=True)
        txt(sl, desc,  4.8, y, 7.8, 0.32, size=12, colour=MUTED)
        y += 0.36

    # ── 12. VOICE ─────────────────────────────────────────────────────────
    sl = add_slide(prs)
    fill_bg(sl)
    box(sl, 0, 0, 0.08, 7.5, bg=RGBColor(0xec, 0x48, 0x99))
    txt(sl, "Card 9 — Voice Call", 0.4, 0.25, 12, 0.6, size=30, bold=True)
    tag_pill(sl, "Voice", 0.4, 0.9, RGBColor(0xec, 0x48, 0x99))
    divider(sl, 1.3)

    txt(sl, "Outbound Twilio discovery call — real, not a mock.",
        0.4, 1.45, 12, 0.45, size=18, bold=True, colour=WHITE)

    details = [
        ("Call SID",  "CA4eb61f333e8635f6213b5b3c980e3ee3"),
        ("Status",    "queued"),
        ("To",        "+251963055269"),
        ("Webhook",   "/webhooks/voice  →  TwiML greeting + IVR menu"),
    ]
    y = 2.1
    for label, val in details:
        box(sl, 0.4, y, 12.3, 0.42, bg=SURFACE,
            border=RGBColor(0x25, 0x2d, 0x45))
        txt(sl, label, 0.55, y+0.06, 2.0, 0.32, size=12, colour=MUTED, bold=True)
        txt(sl, val,   2.7,  y+0.06, 9.8, 0.32, size=12, colour=WHITE)
        y += 0.48

    txt(sl, "Channel Priority", 0.4, 3.6, 12, 0.38, size=14, bold=True, colour=WHITE)
    prio = [
        ("Primary",   "Email  — signal-grounded outreach",               BLUE),
        ("Secondary", "SMS    — automated handoff on warm reply",         GREEN),
        ("Top tier",  "Voice  — high-value confirmed leads only",         RGBColor(0xec, 0x48, 0x99)),
    ]
    y = 4.05
    for tier, desc, col in prio:
        txt(sl, tier, 0.4, y, 1.6, 0.35, size=13, bold=True, colour=col)
        txt(sl, desc, 2.1, y, 10,  0.35, size=13, colour=WHITE)
        y += 0.45

    # ── 13. COST + EVIDENCE ───────────────────────────────────────────────
    sl = add_slide(prs)
    fill_bg(sl)
    box(sl, 0, 0, 0.08, 7.5, bg=GREEN)
    txt(sl, "Cost Summary + Evidence Graph", 0.4, 0.25, 12, 0.6,
        size=30, bold=True)
    divider(sl, 1.0)

    # cost table
    txt(sl, "invoice_summary.json", 0.4, 1.1, 6, 0.35, size=13, bold=True, colour=AMBER)
    cost_rows = [
        ("τ²-Bench full run (30 tasks × 5 trials)", "gpt-4o-mini",          "$0.885"),
        ("Email compose — 50 runs",                 "gpt-4o-mini",          "$0.011"),
        ("Resend outbound (happy path)",            "Resend API",           "$0.007"),
    ]
    y = 1.52
    for item, model, cost in cost_rows:
        txt(sl, item,  0.4, y, 6.5, 0.32, size=12, colour=WHITE)
        txt(sl, model, 7.0, y, 3.0, 0.32, size=12, colour=MUTED)
        txt(sl, cost,  10.2,y, 1.5, 0.32, size=12, colour=GREEN, bold=True)
        y += 0.36
    divider(sl, y + 0.05)
    txt(sl, "Total spend",                   0.4,  y+0.12, 6.5, 0.32, size=12, bold=True, colour=WHITE)
    txt(sl, "$0.903  ·  within $5.00 target",10.2, y+0.12, 2.8, 0.32, size=12, bold=True, colour=GREEN)

    # evidence graph
    divider(sl, y + 0.6)
    txt(sl, "evidence_graph.json — every PDF number is traced",
        0.4, y+0.72, 12, 0.35, size=13, bold=True, colour=LAVENDER)
    claims = [
        ("C001", "τ²-Bench pass@1 = 72.67%",                    "eval/score_log.json"),
        ("C002", "p50 latency = 2.8s",                          "eval/latency_results.json"),
        ("C005", "Cost per prospect = $0.0059",                  "invoice_summary.json"),
        ("C006", "Happy path trace — email + HubSpot + Cal.com", "Langfuse trace 908c8ce3"),
        ("C007", "All 6 signals present for NovaPay",            "data/hiring_signal_brief_novapay_v2.json"),
    ]
    cy = y + 1.1
    for cid, claim, src in claims:
        txt(sl, cid,   0.4, cy, 0.8, 0.3, size=11, bold=True, colour=LAVENDER)
        txt(sl, claim, 1.3, cy, 7.5, 0.3, size=11, colour=WHITE)
        txt(sl, src,   9.0, cy, 4.1, 0.3, size=10, colour=MUTED)
        cy += 0.33

    # ── 14. KILL SWITCH ───────────────────────────────────────────────────
    sl = add_slide(prs)
    fill_bg(sl)
    box(sl, 0, 0, 0.08, 7.5, bg=RED)
    txt(sl, "Kill-Switch Thresholds", 0.4, 0.25, 12, 0.6, size=30, bold=True)
    divider(sl, 1.0)
    txt(sl, "If any metric breaches its threshold on two consecutive days — all outbound halts.",
        0.4, 1.1, 12, 0.4, size=14, colour=MUTED)

    headers = ["Metric", "Threshold", "Measurement"]
    hx = [0.4, 5.0, 8.5]
    for h, x in zip(headers, hx):
        txt(sl, h, x, 1.62, 3.5, 0.35, size=12, bold=True, colour=MUTED)
    divider(sl, 2.05)

    rows = [
        ("hallucination_rate",        "> 2%",   "LLM-as-judge on random 5% of outbound"),
        ("cost_per_qualified_lead",   "> $8.00","invoice_summary.json ÷ qualified leads"),
        ("icp_conflict_flag_rate",    "> 15%",  "Fraction of prospects with conflict_flag = true"),
        ("opt_out_rate (SMS/email)",  "> 5%",   "STOP commands ÷ total SMS conversations"),
    ]
    y = 2.15
    for metric, threshold, measure in rows:
        bg_c = RGBColor(0x1c, 0x21, 0x38) if rows.index((metric,threshold,measure)) % 2 == 0 else SURFACE
        box(sl, 0.4, y, 12.3, 0.48, bg=bg_c)
        txt(sl, metric,    0.55, y+0.08, 4.3, 0.35, size=13, colour=WHITE)
        txt(sl, threshold, 5.0,  y+0.08, 3.3, 0.35, size=13, bold=True, colour=RED)
        txt(sl, measure,   8.5,  y+0.08, 4.1, 0.35, size=12, colour=MUTED)
        y += 0.54

    txt(sl, "Rollback: two consecutive days above any threshold → routes to staff sink → delivery lead review required.",
        0.4, 5.6, 12.3, 0.5, size=12, colour=AMBER)

    # ── 15. CLOSING ───────────────────────────────────────────────────────
    sl = add_slide(prs)
    fill_bg(sl)
    box(sl, 0, 0, 0.08, 7.5, bg=BLUE)
    box(sl, 0.08, 3.5, 13.25, 0.06, bg=BLUE)
    txt(sl, "The full pipeline.", 0.4, 1.0, 12.5, 0.8,
        size=42, bold=True, colour=WHITE)

    summary = [
        ("6 signals  ·  4 ICP segments  ·  confidence-gated assertion mode", WHITE,    True,  16),
        ("Email  →  SMS  →  Booking  →  HubSpot — zero humans in the loop",  LAVENDER, False, 16),
        ("Every trace logged  ·  Every number in evidence_graph.json",        MUTED,    False, 15),
        ("Zero CFPB fields  ·  Total spend $0.90  ·  Within $5.00 target",   GREEN,    True,  15),
    ]
    multiline(sl, summary, 0.4, 1.95, 12.5, 1.4)

    txt(sl, "Rahel Samson  ·  10Academy Forward-Deployed Challenge  ·  April 2026",
        0.4, 3.7, 12.5, 0.45, size=14, colour=MUTED)

    out_path = "/Users/aman/Desktop/projects/10academy/The-Forward-Deployed-Choice/Tenacious_Consulting_Presentation.pptx"
    prs.save(out_path)
    print(f"Saved: {out_path}")
    return out_path


if __name__ == "__main__":
    build()
